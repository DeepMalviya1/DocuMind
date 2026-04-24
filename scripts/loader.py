"""
loader.py

Simple file loader.
Reads text from different file types and returns it as a string.
Uses vision model for images and graphs inside documents.
"""

import os
import csv
import sqlite3
import tempfile

import pytesseract
import pandas as pd
import fitz
from PIL import Image
from docx import Document
from pptx import Presentation
from odf.opendocument import load as odf_load
from odf.text import P

from scripts.log import get_logger
from scripts.vision import describe_image

class DocumentLoader:

    SUPPORTED = (
        ".txt", ".csv", ".docx",
        ".xls", ".xlsx", ".pptx",
        ".png", ".jpg", ".jpeg", ".bmp", ".tiff",
        ".db", ".pdf", ".odt"
    )

    def __init__(self):
        self.logger = get_logger("DocumentLoader")
        self.logger.info("DocumentLoader ready")

    def load_directory(self, dir_path):
        """Load all supported files in a directory and subfolders."""

        if not os.path.exists(dir_path):
            self.logger.error(f"Directory not found: {dir_path}")
            raise FileNotFoundError(f"Directory not found: {dir_path}")

        if not os.path.isdir(dir_path):
            self.logger.error(f"Not a directory: {dir_path}")
            raise ValueError(f"Not a directory: {dir_path}")

        self.logger.info(f"Scanning directory: {dir_path}")

        results = {}
        failed = {}

        for root, dirs, files in os.walk(dir_path):
            for name in sorted(files):
                file_path = os.path.join(root, name)
                _, ext = os.path.splitext(name)

                if ext.lower() not in self.SUPPORTED:
                    self.logger.debug(f"Skipping unsupported: {file_path}")
                    continue

                try:
                    content = self.load(file_path)
                    results[file_path] = content
                except Exception as e:
                    self.logger.error(f"Failed: {file_path} - {e}")
                    failed[file_path] = str(e)

        self.logger.info(f"Done. Success: {len(results)}, Failed: {len(failed)}")
        return results

    def load(self, file_path):
        """Load a single file and return its text content."""

        if not os.path.exists(file_path):
            self.logger.error(f"File not found: {file_path}")
            raise FileNotFoundError(f"File not found: {file_path}")

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        self.logger.info(f"Loading: {file_path}")

        if ext == ".txt":
            content = self.load_txt(file_path)

        elif ext == ".csv":
            content = self.load_csv(file_path)

        elif ext == ".docx":
            content = self.load_docx(file_path)

        elif ext in (".xls", ".xlsx"):
            content = self.load_excel(file_path)

        elif ext == ".pptx":
            content = self.load_pptx(file_path)

        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
            content = self.load_image(file_path)

        elif ext == ".db":
            content = self.load_db(file_path)

        elif ext == ".pdf":
            content = self.load_pdf(file_path)

        elif ext == ".odt":
            content = self.load_odt(file_path)

        else:
            self.logger.error(f"Unsupported file type: {ext}")
            raise ValueError(f"Unsupported file type: {ext}")

        self.logger.info(f"Got {len(content)} characters from {file_path}")
        return content

    # ------------------------------------------------------------------
    # Text
    # ------------------------------------------------------------------
    def load_txt(self, file_path):
        self.logger.debug("Reading text file")
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # ------------------------------------------------------------------
    # CSV
    # ------------------------------------------------------------------
    def load_csv(self, file_path):
        self.logger.debug("Reading csv file")
        lines = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            for row in csv.reader(f):
                lines.append(" | ".join(row))
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # Word (.docx) - extracts text, tables and images
    # ------------------------------------------------------------------
    def load_docx(self, file_path):
        self.logger.debug("Reading docx file")
        doc = Document(file_path)
        lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)

        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))

        # Extract and describe images
        image_descriptions = self._extract_docx_images(doc)
        if image_descriptions:
            lines.append("")
            lines.append("--- Images found in document ---")
            for i, desc in enumerate(image_descriptions, start=1):
                lines.append(f"[Image {i}]: {desc}")

        return "\n".join(lines)

    def _extract_docx_images(self, doc):
        """Extract images from a docx file and describe them."""

        descriptions = []

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    temp_path = self._save_temp_image(image_data)
                    desc = describe_image(temp_path)
                    if desc:
                        descriptions.append(desc)
                    os.remove(temp_path)
                except Exception as e:
                    self.logger.error(f"Failed to process docx image: {e}")

        return descriptions

    # ------------------------------------------------------------------
    # Excel (.xls, .xlsx)
    # ------------------------------------------------------------------
    def load_excel(self, file_path):
        self.logger.debug("Reading excel file")
        sheets = pd.read_excel(file_path, sheet_name=None)
        lines = []

        for name, df in sheets.items():
            lines.append(f"--- Sheet: {name} ---")
            lines.append(df.to_string(index=False))
            lines.append("")

        return "\n".join(lines)

    # ------------------------------------------------------------------
    # PowerPoint (.pptx) - extracts text, tables and images
    # ------------------------------------------------------------------
    def load_pptx(self, file_path):
        self.logger.debug("Reading pptx file")
        prs = Presentation(file_path)
        lines = []

        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {i} ---")

            image_count = 0

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        lines.append(" | ".join(cells))

                if shape.shape_type == 13:
                    image_count += 1
                    try:
                        image_data = shape.image.blob
                        temp_path = self._save_temp_image(image_data)
                        desc = describe_image(temp_path)
                        if desc:
                            lines.append(f"[Image {image_count}]: {desc}")
                        os.remove(temp_path)
                    except Exception as e:
                        self.logger.error(f"Failed to process pptx image: {e}")

            lines.append("")

        return "\n".join(lines)

    # ------------------------------------------------------------------
    # Images (.png, .jpg, etc) - uses Groq vision
    # ------------------------------------------------------------------
    def load_image(self, file_path):
        self.logger.debug("Reading image")
        return describe_image(file_path)

    # ------------------------------------------------------------------
    # PDF (.pdf) - extracts text, OCR for scanned pages, images
    # ------------------------------------------------------------------
    def load_pdf(self, file_path):
        self.logger.debug("Reading pdf file")
        doc = fitz.open(file_path)
        lines = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text
            text = page.get_text().strip()

            # If no text, try OCR
            if not text:
                self.logger.debug(f"Page {page_num + 1}: no text, trying OCR")
                text = self._ocr_pdf_page(page)

            if text:
                lines.append(f"--- Page {page_num + 1} ---")
                lines.append(text)

            # Extract images from the page
            image_descriptions = self._extract_pdf_images(page, page_num)
            if image_descriptions:
                for desc in image_descriptions:
                    lines.append(desc)

            lines.append("")

        doc.close()
        return "\n".join(lines)

    def _ocr_pdf_page(self, page):
        """Convert a PDF page to image and run OCR."""

        zoom = 300 / 72
        matrix = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=matrix)

        image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        text = pytesseract.image_to_string(image)

        return text.strip()

    def _extract_pdf_images(self, page, page_num):
        """Extract images from a PDF page and describe them."""

        descriptions = []
        image_list = page.get_images(full=True)

        if not image_list:
            return descriptions

        self.logger.debug(f"Page {page_num + 1}: found {len(image_list)} images")

        doc = page.parent

        for img_index, img_info in enumerate(image_list, start=1):
            try:
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                image_data = base_image["image"]

                # Skip very small images (icons, bullets)
                if len(image_data) < 5000:
                    continue

                temp_path = self._save_temp_image(image_data)
                desc = describe_image(temp_path)

                if desc:
                    descriptions.append(
                        f"[Page {page_num + 1}, Image {img_index}]: {desc}"
                    )

                os.remove(temp_path)

            except Exception as e:
                self.logger.error(f"Failed to process PDF image: {e}")

        return descriptions

    # ------------------------------------------------------------------
    # OpenDocument Text (.odt)
    # ------------------------------------------------------------------
    def load_odt(self, file_path):
        self.logger.debug("Reading odt file")
        doc = odf_load(file_path)
        lines = []

        paragraphs = doc.getElementsByType(P)

        for para in paragraphs:
            text = ""
            for node in para.childNodes:
                if hasattr(node, "data"):
                    text += node.data
                elif hasattr(node, "__str__"):
                    text += str(node)

            text = text.strip()
            if text:
                lines.append(text)

        return "\n".join(lines)

    # ------------------------------------------------------------------
    # SQLite Database (.db)
    # ------------------------------------------------------------------
    # def load_db(self, file_path):
    #     self.logger.debug("Reading database")
    #     conn = sqlite3.connect(file_path)
    #     cursor = conn.cursor()
    #     lines = []

    #     cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    #     tables = cursor.fetchall()

    #     if not tables:
    #         conn.close()
    #         return "Database is empty."

    #     for (table_name,) in tables:
    #         lines.append(f"--- Table: {table_name} ---")

    #         cursor.execute(f"PRAGMA table_info('{table_name}');")
    #         columns = [col[1] for col in cursor.fetchall()]
    #         lines.append("Columns: " + " | ".join(columns))

    #         cursor.execute(f"SELECT * FROM '{table_name}';")
    #         rows = cursor.fetchall()

    #         if not rows:
    #             lines.append("(no data)")
    #         else:
    #             for row in rows:
    #                 lines.append(" | ".join(str(v) for v in row))

    #         lines.append("")

    #     conn.close()
    #     return "\n".join(lines)
        # ------------------------------------------------------------------
    # SQLite Database (.db) - with metadata
    # ------------------------------------------------------------------
    def load_db(self, file_path):
        self.logger.debug("Reading database")
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        lines = []

        # Get all tables
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()

        if not tables:
            conn.close()
            return "Database is empty."

        file_name = os.path.basename(file_path)
        lines.append(f"Database: {file_name}")
        lines.append(f"Total tables: {len(tables)}")
        lines.append("")

        for (table_name,) in tables:

            # Get column info
            cursor.execute(f"PRAGMA table_info('{table_name}');")
            columns_info = cursor.fetchall()
            column_names = [col[1] for col in columns_info]

            # Get row count
            cursor.execute(f"SELECT COUNT(*) FROM '{table_name}';")
            row_count = cursor.fetchone()[0]

            # Table info
            lines.append(f"--- Table: {table_name} ---")
            lines.append(f"Rows: {row_count}, Columns: {len(column_names)}")

            # Column names with types
            for col in columns_info:
                col_name = col[1]
                col_type = col[2] if col[2] else "UNKNOWN"
                lines.append(f"  - {col_name} ({col_type})")

            # Simple stats for number columns
            for col in columns_info:
                col_name = col[1]
                col_type = (col[2] or "").upper()

                if col_type in ("INTEGER", "REAL", "NUMERIC", "FLOAT", "INT"):
                    try:
                        cursor.execute(f"""
                            SELECT MIN("{col_name}"), MAX("{col_name}"), 
                                   AVG("{col_name}"), SUM("{col_name}")
                            FROM '{table_name}'
                        """)
                        result = cursor.fetchone()
                        if result[0] is not None:
                            avg_val = round(result[2], 2)
                            lines.append(
                                f"  {col_name} stats: "
                                f"min={result[0]}, max={result[1]}, "
                                f"avg={avg_val}, sum={result[3]}"
                            )
                    except Exception:
                        pass

            # Data
            lines.append("")
            lines.append(" | ".join(column_names))

            cursor.execute(f"SELECT * FROM '{table_name}';")
            rows = cursor.fetchall()

            if not rows:
                lines.append("(no data)")
            else:
                for row in rows:
                    lines.append(" | ".join(str(v) for v in row))

            lines.append("")

        conn.close()
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # Helper
    # ------------------------------------------------------------------
    def _save_temp_image(self, image_data):
        """Save image bytes to a temporary file. Returns the file path."""

        temp_file = tempfile.NamedTemporaryFile(
            suffix=".png", delete=False
        )
        temp_file.write(image_data)
        temp_file.close()

        return temp_file.name